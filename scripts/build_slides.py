"""Genera la presentación de defensa del TFG (PowerPoint). Salida: defensa_tfg.pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

AZUL       = RGBColor(0x1F, 0x3E, 0x6E)
AZUL_CLARO = RGBColor(0x2E, 0x6D, 0xB8)
GRIS       = RGBColor(0x4A, 0x4A, 0x4A)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
VERDE      = RGBColor(0x27, 0x9E, 0x6A)
ROJO       = RGBColor(0xC0, 0x39, 0x2B)
AMARILLO   = RGBColor(0xF3, 0x9C, 0x12)

W = Inches(13.33)  # 16:9 widescreen
H = Inches(7.5)


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=GRIS, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_slide(prs, title_text, bullets, notes=""):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    tf = slide.shapes.title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = AZUL

    content = slide.placeholders[1]
    tf2 = content.text_frame
    tf2.clear()
    tf2.word_wrap = True

    for i, (bullet, level, color) in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.level = level
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(18) if level == 0 else Pt(15)
        run.font.color.rgb = color
        run.font.bold = (level == 0)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def set_slide_bg(slide, color):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, AZUL)

    left, top, width, height = 0, Inches(5.5), W, Inches(2.0)
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BLANCO
    rect.line.fill.background()

    add_textbox(slide, title,
                Inches(0.6), Inches(1.2), Inches(12), Inches(1.8),
                font_size=36, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)
    add_textbox(slide, subtitle,
                Inches(0.6), Inches(3.0), Inches(12), Inches(1.0),
                font_size=20, bold=False, color=AZUL_CLARO, align=PP_ALIGN.LEFT,
                italic=True)
    add_textbox(slide, meta,
                Inches(0.6), Inches(5.6), Inches(12), Inches(1.6),
                font_size=15, bold=False, color=GRIS, align=PP_ALIGN.LEFT)
    return slide


def add_section_divider(prs, section_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, AZUL_CLARO)
    add_textbox(slide, section_text,
                Inches(1.0), Inches(2.8), Inches(11), Inches(2.0),
                font_size=36, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    add_title_slide(
        prs,
        title    = "LoRA Modal para Control de Modos Musicales en ACE-Step 1.5",
        subtitle = "Fine-tuning eficiente para condicionamiento armónico en modelos generativos de música",
        meta     = (
            "Ignacio Franco Almendárez\n"
            "Tutor: Francisco Serradilla\n"
            "Grado en Ciencia de Datos e IA · UPM ETSISI · 2026"
        ),
    )

    add_bullet_slide(prs,
        title_text="¿Por qué controlar la modalidad musical?",
        bullets=[
            ("Modelos actuales generan audio convincente, pero el control armónico de alto nivel es limitado", 0, GRIS),
            ("Puedes pedir 'jazz melancólico'… pero no garantizar que use modo dórico en lugar de menor natural", 1, GRIS),
            ("La modalidad define el carácter emocional de la música: dórico ≠ frigio ≠ eólico", 1, GRIS),
            ("Aplicaciones reales: composición asistida, bandas sonoras, prototipado armónico", 0, GRIS),
            ("Estrategia elegida: LoRA sobre ACE-Step 1.5", 0, AZUL_CLARO),
            ("Fine-tuning eficiente (<1 % de parámetros), sin reentrenar el modelo base (≈1 000 M params)", 1, GRIS),
            ("Permite iterar en hardware universitario con coste razonable", 1, GRIS),
        ],
        notes=(
            "Explicar que la modalidad es una propiedad de alto nivel: no basta timbre o textura, "
            "hay que mantener coherencia armónica global durante toda la pieza. "
            "Mencionar que ACE-Step 1.5 es el modelo base open-source más avanzado disponible."
        ),
    )

    add_bullet_slide(prs,
        title_text="Objetivos y preguntas de investigación",
        bullets=[
            ("Objetivo general: incorporar control modal explícito preservando calidad musical", 0, AZUL),
            ("RQ1  ¿Puede LoRA mejorar el control modal sin degradar calidad global?", 0, GRIS),
            ("RQ2  ¿Qué configuración de hiperparámetros da mejor equilibrio?", 0, GRIS),
            ("RQ3  ¿Cómo evaluar modalidad de forma reproducible?", 0, GRIS),
            ("Alcance: experimental y acotado", 0, AZUL),
            ("No se reentrana el modelo base · Dataset propio · Evaluación dual automática + manual", 1, GRIS),
        ],
        notes="Leer las 3 RQs despacio. Anticipar que las respuestas aparecerán en resultados.",
    )

    add_bullet_slide(prs,
        title_text="¿Qué es un modo musical?",
        bullets=[
            ("Los 7 modos diatónicos comparten las mismas 7 notas — lo que cambia es el punto de partida", 0, AZUL),
            ("Mismo conjunto de notas, distinto centro de gravedad → distinto patrón de tonos y semitonos", 1, GRIS),
            ("Eso cambia qué intervalos caen sobre la tónica… y por tanto el carácter emocional", 1, GRIS),
            ("Ejemplo: D menor natural (eólico) vs D dórico", 0, AZUL),
            ("Eólico:   D · E · F · G · A · Bb · C    →  6ª menor (Sib) →  oscuro, clásico", 1, ROJO),
            ("Dórico:   D · E · F · G · A · B♮ · C    →  6ª mayor (Si♮) →  melancolía jazz, esperanza", 1, VERDE),
            ("Una sola nota diferente → carácter completamente distinto", 1, AMARILLO),
            ("El reto: enseñar al modelo a mantener esa nota característica durante toda la pieza", 0, AZUL),
            ("ACE-Step 1.5 domina la sonoridad mayor pero tiende a ignorar las notas características de los modos menores", 1, GRIS),
        ],
        notes=(
            "Hacer pausa aquí. Tocar el ejemplo si es posible (piano virtual o simplemente cantarlo). "
            "El Si natural vs Si bemol en D dórico es el ejemplo perfecto porque la diferencia es mínima "
            "pero el efecto perceptivo es enorme. "
            "Mencionar que Dorian es el modo de 'So What' de Miles Davis."
        ),
    )

    add_bullet_slide(prs,
        title_text="Marco técnico: ACE-Step 1.5 y LoRA",
        bullets=[
            ("ACE-Step 1.5: modelo generativo musical open-source basado en Flow Matching + DiT", 0, AZUL),
            ("Encoder textual Qwen3-0.6B  →  DiT (bloques Cross-Attention)  →  VAE 1D → audio 48 kHz", 1, GRIS),
            ("LoRA (Low-Rank Adaptation)", 0, AZUL),
            ("W' = W + (α/r)·BA  —  solo A y B reciben gradiente", 1, GRIS),
            ("Módulos objetivo: q_proj, k_proj, v_proj, o_proj de Cross-Attention del DiT", 1, GRIS),
            ("<1 % del total de parámetros entrenables · Adaptadores ligeros (~MB vs ~GB)", 1, VERDE),
            ("Función de pérdida: Flow Matching  E[‖v_θ(z_t, t, c) − (z₁ − z₀)‖²]", 0, GRIS),
        ],
        notes=(
            "Mostrar figura de arquitectura si se puede. "
            "Insistir: todo el modelo base queda congelado, solo actualizamos A y B. "
            "La hipótesis es que la Cross-Attention aprende la señal modal adicional."
        ),
    )

    add_bullet_slide(prs,
        title_text="Dataset y pipeline experimental",
        bullets=[
            ("Dataset propio: 105 clips de 20–50 s, 7 modos diatónicos, etiquetado manual", 0, AZUL),
            ("Jónico: 7  ·  Dórico: 21  ·  Frigio: 16  ·  Lidio: 18  ·  Mixolidio: 18  ·  Eólico: 14  ·  Locrio: 8", 1, GRIS),
            ("Desbalance asumido: Jónico infrarepresentado, Dórico sobrerepresentado", 1, AMARILLO),
            ("Pipeline completo", 0, AZUL),
            ("Audio → VAE (congelado) → tensores  →  entrenamiento LoRA  →  inferencia por lotes  →  métricas automáticas  →  evaluación manual", 1, GRIS),
            ("Evaluación en dos fases (reducción de carga sin perder rigor)", 0, AZUL),
            ("Fase 1: auto_score_0_1 para cribar checkpoints (cromagrama + plantillas Krumhansl–Kessler)", 1, GRIS),
            ("Fase 2: rúbrica manual 4 criterios × escala 1–5: centro tonal, color modal, no modulación, coherencia", 1, GRIS),
        ],
        notes=(
            "Remarcar que los audios no se redistribuyen, solo los tensores VAE (no invertibles). "
            "La evaluación manual fue realizada por el autor: limitación declarada, un solo evaluador."
        ),
    )

    add_bullet_slide(prs,
        title_text="Búsqueda sistemática: 24 configuraciones",
        bullets=[
            ("Grid: rank ∈ {32, 64, 128}  ×  alpha ∈ {32, 128}  ×  lr ∈ {10⁻⁴, 10⁻³}  ×  dropout ∈ {0.05, 0.10}", 0, GRIS),
            ("Evaluación uniforme: 7 modos × 2 semillas × checkpoint epoch 140", 1, GRIS),
            ("Resultado principal: la tasa de aprendizaje es el hiperparámetro dominante", 0, AZUL),
            ("lr = 10⁻⁴  →  media 0.7217  (σ = 0.0029)", 1, VERDE),
            ("lr = 10⁻³  →  media 0.6932  (σ = 0.0103)  —  3.5× más varianza", 1, ROJO),
            ("Δ = 0.029 entre grupos de lr, frente a Δ < 0.01 de rank y dropout", 1, GRIS),
            ("Rank y dropout son prácticamente irrelevantes en este contexto", 0, AMARILLO),
            ("Configuración dominante: r=32, α=32, lr=10⁻⁴  →  validada con experimento de confirmación", 1, AZUL_CLARO),
        ],
        notes=(
            "El heatmap y el gráfico de sensibilidad ilustran bien este resultado. "
            "La mayor varianza con lr=1e-3 significa que el entrenamiento es inestable, "
            "no que el resultado final sea peor de media."
        ),
    )

    add_bullet_slide(prs,
        title_text="Hallazgo metodológico: métrica auto ≠ escucha experta",
        bullets=[
            ("Comparativa dentro del top-5 de la búsqueda sistemática", 0, AZUL),
            ("Exp_01 (r32, α32, lr=10⁻⁴, do=0.05)  →  auto_score = 0.7244  —  1º en ranking", 1, GRIS),
            ("                                          →  manual = 8.48 / 20  —  el PEOR del top-5", 1, ROJO),
            ("Exp_18 (r128, α32, lr=10⁻⁴, do=0.10)  →  auto_score = 0.7241  —  4º en ranking", 1, GRIS),
            ("                                           →  manual = 9.89 / 20  —  el MEJOR del top-5", 1, VERDE),
            ("Diferencia de 0.0003 en auto_score oculta 1.4 puntos en calidad auditiva real", 0, AMARILLO),
            ("Conclusión: la métrica automática es útil para cribar, no para decidir", 0, AZUL),
            ("→ Justifica el protocolo de evaluación dual adoptado (responde RQ3)", 1, GRIS),
        ],
        notes=(
            "Este slide responde directamente RQ3. "
            "El auto_score es bueno para descartar configuraciones malas, "
            "pero no discrimina entre las buenas."
        ),
    )

    add_bullet_slide(prs,
        title_text="Patrón de confusión: sesgo hacia sonoridades mayores",
        bullets=[
            ("El modelo base de ACE-Step tiene un sesgo inherente hacia resoluciones en mayor", 0, ROJO),
            ("Dórico (D menor)   → genera D mayor de forma sistemática", 1, GRIS),
            ("Frigio (E menor)   → genera E mayor, no capta la 2ª disminuida", 1, GRIS),
            ("Mixolidio (G mix.) → genera G mayor natural, omite la 7ª menor", 1, GRIS),
            ("Locrio             → inestable, resuelve en B/E mayor", 1, GRIS),
            ("LoRA reduce pero no elimina este sesgo", 0, AMARILLO),
            ("Jónico y Lidio funcionan muy bien: el modelo base ya los domina", 1, VERDE),
            ("Implicación práctica: mejorar modos menores requiere más datos o señal más explícita", 1, GRIS),
        ],
        notes=(
            "Aquí se puede mencionar que Jónico tiene solo 7 clips de entrenamiento "
            "y aun así es el mejor modo: es el modelo base quien ya lo domina. "
            "La LoRA apenas necesita 'recordárselo'."
        ),
    )

    add_bullet_slide(prs,
        title_text="Resultados del checkpoint de referencia (run_D, epoch 100)",
        bullets=[
            ("Configuración: r=64, α=128, lr=5×10⁻⁵, dropout=0.05  —  mejor de la primera tanda", 0, AZUL),
            ("Media global: 12.41 / 20  (62.1 % del máximo)  ·  rúbrica 4 criterios × 1–5", 0, GRIS),
            ("Jónico: 17.75  ·  Lidio: 15.50  ·  Mixolidio: 12.00", 1, VERDE),
            ("Dórico: 11.75  ·  Eólico: 11.38  ·  Frigio: 11.00  ·  Locrio: 7.50", 1, ROJO),
            ("Prompt mínimo ('modo solo') supera al detallado en promedio global y en dórico", 0, AMARILLO),
            ("Demasiado contexto textual puede interferir con la señal modal del adaptador", 1, GRIS),
            ("La búsqueda de hiperparámetros identifica una config que mejora los modos menores → slide siguiente", 0, AZUL_CLARO),
        ],
        notes=(
            "Señalar el gradiente claro: modos mayores bien, modos menores y locrio mal. "
            "El dato del prompt mínimo es sorprendente y responde parcialmente RQ2."
        ),
    )

    add_bullet_slide(prs,
        title_text="Experimento de confirmación: validación del ciclo completo",
        bullets=[
            ("Config óptima de la búsqueda: r=32, α=32, lr=10⁻⁴, dropout=0.05  (epoch 140)", 0, AZUL),
            ("Global: 11.14 / 20  (55.7 %)  vs  run_D: 12.41 / 20  (62.1 %) — globalmente inferior", 0, AMARILLO),
            ("Objetivo cumplido: reducir el sesgo en modos menores", 0, AZUL),
            ("Dórico:  run_D → 11.75/20 (58.8 %)  ·  Confirm → 14.13/20 (70.7 %)   +12 pp", 1, VERDE),
            ("Eólico:  run_D → 11.38 (56.9 %)  ·  Confirm → 10.00 (50.0 %)", 1, GRIS),
            ("Primera config con ejemplos dóricos reconocibles de forma consistente", 1, GRIS),
            ("6ª mayor prominente, variaciones armónicas características del modo dórico", 1, GRIS),
            ("La métrica automática va en sentido contrario: confirm 0.705 vs run_D 0.711", 0, ROJO),
            ("→ Tercer caso donde auto_score infravalora una config perceptivamente superior en modos menores", 1, GRIS),
            ("Ciclo demostrado: búsqueda → config óptima → confirmación empírica focalizada", 0, AZUL_CLARO),
        ],
        notes=(
            "Ser honesto: globalmente la config de confirmación es ligeramente peor que run_D. "
            "El valor está en el objetivo específico: reducir el sesgo en dórico. "
            "La LoRA desplaza capacidad del modelo base hacia modos menores, "
            "lo que puede sacrificar algo de rendimiento en modos mayores ya dominados."
        ),
    )

    add_bullet_slide(prs,
        title_text="Conclusiones",
        bullets=[
            ("RQ1 ✓  LoRA sobre ACE-Step 1.5 permite control modal sin reentrenamiento completo", 0, VERDE),
            ("Calidad general preservada · Coste computacional asumible · Adaptadores ligeros", 1, GRIS),
            ("RQ2 ✓  lr = 10⁻⁴ y α = 32 son la configuración dominante (Δ = 0.029)", 0, VERDE),
            ("Confirm: global 56 % vs 62 % run_D · Pero dórico 59 % → 71 % (+12 pp): primer resultado consistente", 1, GRIS),
            ("RQ3 ✓  Evaluación dual necesaria: auto_score criba, escucha experta decide", 0, VERDE),
            ("Diferencias de 0.0003 en auto_score pueden ocultar 1.4 puntos de calidad perceptiva", 1, GRIS),
            ("Limitación principal: sesgo del modelo base hacia mayor, un solo evaluador, n=2 semillas", 0, AMARILLO),
            ("Aportación clave: protocolo experimental reproducible para control armónico con LoRA", 0, AZUL),
        ],
        notes="Recapitular las 3 RQs con sus respuestas. Ir directo, sin rodeos.",
    )

    add_bullet_slide(prs,
        title_text="Trabajo futuro y cierre",
        bullets=[
            ("Líneas de continuación", 0, AZUL),
            ("Ampliar dataset en modos menores · Sobrerepresentar dórico y frigio", 1, GRIS),
            ("Evaluar checkpoints intermedios (no solo epoch final) para localizar óptimo con más precisión", 1, GRIS),
            ("Incorporar métricas automáticas más sensibles al color modal fino", 1, GRIS),
            ("Replicar con múltiples evaluadores para obtener acuerdo inter-evaluadores", 1, GRIS),
            ("Aportación principal de este trabajo", 0, AZUL),
            ("Un procedimiento experimental replicable para estudiar control armónico en modelos generativos musicales con recursos acotados", 1, AZUL_CLARO),
            ("Código disponible en el repositorio del TFG", 1, GRIS),
        ],
        notes="Agradecer y abrir turno de preguntas.",
    )

    import sys
    out = sys.argv[1] if len(sys.argv) > 1 else "defensa_tfg.pptx"
    prs.save(out)
    print(f"Presentación guardada en: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
